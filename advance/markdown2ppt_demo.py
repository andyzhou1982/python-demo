from pptx import Presentation
from pptx.util import Inches
import re
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

def markdown_to_slides(md_content, template_pptx):
    # 使用正则表达式来解析Markdown中的标题和内容
    pattern = re.compile(r'^#+ (.*)', re.MULTILINE)
    matches = pattern.findall(md_content)

    # 加载现有的PPT模板
    presentation = Presentation(template_pptx)

    # 解析Markdown内容并添加到PPT中
    for i, match in enumerate(matches, start=1):
        slide_layout_index = 0  # 使用模板的第一张幻灯片布局
        slide = presentation.slides.add_slide(slide_layout_index)

        # 添加标题
        title = slide.shapes.title
        title.text = match

        # 如果有多级标题，解析子标题
        sub_pattern = re.compile(r'###+ (.*)', re.MULTILINE)
        sub_matches = sub_pattern.findall(md_content)
        
        if sub_matches:
            # 添加子标题
            sub_title = slide.placeholders[1]
            sub_title.text = sub_matches[i-1]

        # 添加内容
        content_pattern = re.compile(r'^>(.*)', re.MULTILINE)
        content_matches = content_pattern.findall(md_content)
        
        if content_matches:
            # 添加到幻灯片内容区域
            slide.placeholders[3].text = content_matches[i-1]

    # 保存PPT文件
    presentation.save('markdown_to_pptx.pptx')


def generate(template_pptx):
    prs = Presentation()
    slide_layouts = prs.slide_layouts
    slide = prs.slides.add_slide(slide_layouts[0])
    title = slide.placeholders[0]
    subtitle = slide.placeholders[1]
    title.text = "新能源汽车年终总结报告"
    subtitle.text = "驶向未来：新能源汽车行业年终总结"
    prs.slides.add_slide(slide_layouts[1])
    prs.slides.add_slide(slide_layouts[7])
    '''
    end = prs.slides.add_slide(slide_layouts[6])
    text_box = end.shapes.add_textbox(Inches(1),Inches(1),Inches(1),Inches(1))
    text_frame  = text_box.text_frame
    text_frame.text = "谢谢观看"
    text_frame.paragraphs[0].font.size = 36
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].alignment =1
    print(len(end.shapes))
    '''
    prs.save('resources/test.pptx')

def test():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))


# 示例Markdown内容
md_content = """
# 演示标题

这是演示的主标题。

## 子标题

这是子标题及其内容。

> 这里是引用内容，
> 它应当被识别并显示为引用。

## 另一张幻灯片

这里是另一张幻灯片的内容。
"""

# 模板PPT文件路径
template_pptx = 'resources/template1.pptx'

# 调用函数转换Markdown到PPT
#markdown_to_slides(md_content, template_pptx)
generate(template_pptx)

#test()

